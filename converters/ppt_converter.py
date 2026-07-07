"""PPT (pptx) 相关格式转换"""

from pathlib import Path


def ppt_to_pdf(input_path: Path, output_path: Path) -> Path:
    """PPT → PDF"""
    from pptx import Presentation
    from reportlab.lib.pagesizes import landscape, A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    import html as html_module
    import os

    font_name = "Helvetica"
    for fp in ["C:/Windows/Fonts/simhei.ttf", "C:/Windows/Fonts/simsun.ttc"]:
        if os.path.exists(fp):
            try:
                pdfmetrics.registerFont(TTFont("ChineseFont", fp))
                font_name = "ChineseFont"
                break
            except Exception:
                continue

    prs = Presentation(str(input_path))
    page_size = landscape(A4)
    pdf_doc = SimpleDocTemplate(str(output_path), pagesize=page_size)

    title_style = ParagraphStyle(
        "SlideTitle", fontName=font_name, fontSize=24,
        leading=32, spaceAfter=20, spaceBefore=20,
    )
    body_style = ParagraphStyle(
        "SlideBody", fontName=font_name, fontSize=14,
        leading=22, spaceAfter=8,
    )
    page_num_style = ParagraphStyle(
        "PageNum", fontName=font_name, fontSize=10,
        textColor="#999999", alignment=1,
    )

    elements = []
    for i, slide in enumerate(prs.slides):
        if i > 0:
            elements.append(PageBreak())

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                text = html_module.escape(text)
                is_title = (
                    shape.shape_type is not None
                    and hasattr(shape, "placeholder_format")
                    and shape.placeholder_format is not None
                    and shape.placeholder_format.idx == 0
                )
                style = title_style if is_title else body_style
                try:
                    elements.append(Paragraph(text, style))
                except Exception:
                    elements.append(Paragraph(text.replace("&", "&amp;"), style))

        elements.append(Spacer(1, 20))
        elements.append(Paragraph(f"— 第 {i+1} 页 —", page_num_style))

    if not elements:
        elements.append(Spacer(1, 1))

    pdf_doc.build(elements)
    return output_path


def ppt_to_image(input_path: Path, output_path: Path, fmt: str = "png") -> Path:
    """PPT → 图片（简化：将每页内容渲染为图片）"""
    from pptx import Presentation
    from PIL import Image, ImageDraw, ImageFont

    prs = Presentation(str(input_path))
    output_paths = []
    stem = output_path.stem
    parent = output_path.parent
    width, height = 1920, 1080

    for i, slide in enumerate(prs.slides):
        img = Image.new("RGB", (width, height), "white")
        draw = ImageDraw.Draw(img)

        y_offset = 100
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                # 使用系统字体
                try:
                    font = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 28)
                except Exception:
                    font = ImageFont.load_default()

                # 自动换行
                lines = _wrap_text(text, draw, font, width - 200)
                for line in lines:
                    if y_offset > height - 100:
                        break
                    draw.text((100, y_offset), line, fill="#333333", font=font)
                    y_offset += 40
                y_offset += 10

        # 添加页码
        try:
            small_font = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 18)
        except Exception:
            small_font = ImageFont.load_default()
        draw.text((width - 100, height - 50), f"{i+1}/{len(prs.slides)}", fill="#999999", font=small_font)

        page_path = parent / f"{stem}_slide{i+1}.{fmt}"
        img.save(str(page_path))
        output_paths.append(page_path)

    return output_paths[0] if output_paths else output_path


def _wrap_text(text: str, draw, font, max_width: int) -> list[str]:
    """文本自动换行"""
    lines = []
    current_line = ""
    for char in text:
        test_line = current_line + char
        bbox = draw.textbbox((0, 0), test_line, font=font)
        if bbox[2] - bbox[0] > max_width:
            if current_line:
                lines.append(current_line)
            current_line = char
        else:
            current_line = test_line
    if current_line:
        lines.append(current_line)
    return lines or [""]


def ppt_to_html(input_path: Path, output_path: Path) -> Path:
    """PPT → HTML"""
    from pptx import Presentation
    import html as html_module

    prs = Presentation(str(input_path))
    slides_html = []

    for i, slide in enumerate(prs.slides):
        slide_parts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = html_module.escape(para.text.strip())
                if not text:
                    continue
                is_title = (
                    hasattr(shape, "placeholder_format")
                    and shape.placeholder_format is not None
                    and shape.placeholder_format.idx == 0
                )
                if is_title:
                    slide_parts.append(f'<h2 style="color:#1a56db;margin:0 0 15px;">{text}</h2>')
                else:
                    slide_parts.append(f'<p style="margin:5px 0;">{text}</p>')

        slides_html.append(
            f'<div class="slide" style="background:white;padding:40px;margin:20px auto;'
            f'max-width:960px;min-height:500px;box-shadow:0 2px 15px rgba(0,0,0,0.1);'
            f'border-radius:8px;page-break-after:always;">'
            f'<div style="text-align:right;color:#ccc;font-size:12px;">第 {i+1} 页</div>'
            f'{"".join(slide_parts)}'
            f'</div>'
        )

    html_content = f"""<!DOCTYPE html>
<html><head><meta charset="utf-8">
<title>{input_path.stem}</title>
<style>body{{background:#f0f2f5;font-family:system-ui,sans-serif;padding:20px;}}</style>
</head><body>
{''.join(slides_html)}
</body></html>"""
    output_path.write_text(html_content, encoding="utf-8")
    return output_path
