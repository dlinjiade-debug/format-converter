"""Word (docx) 相关格式转换"""

from pathlib import Path


def word_to_pdf(input_path: Path, output_path: Path) -> Path:
    """Word → PDF"""
    try:
        from docx2pdf import convert
        convert(str(input_path), str(output_path))
        return output_path
    except Exception:
        # 回退方案：用 python-docx 读取 + reportlab 生成 PDF
        return _word_to_pdf_fallback(input_path, output_path)


def _word_to_pdf_fallback(input_path: Path, output_path: Path) -> Path:
    """Word → PDF 回退方案"""
    from docx import Document
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    import os

    # 尝试注册中文字体
    font_name = "Helvetica"
    chinese_font_paths = [
        "C:/Windows/Fonts/simhei.ttf",
        "C:/Windows/Fonts/simsun.ttc",
        "C:/Windows/Fonts/msyh.ttc",
    ]
    for fp in chinese_font_paths:
        if os.path.exists(fp):
            try:
                pdfmetrics.registerFont(TTFont("ChineseFont", fp))
                font_name = "ChineseFont"
                break
            except Exception:
                continue

    doc = Document(str(input_path))
    pdf_doc = SimpleDocTemplate(
        str(output_path), pagesize=A4,
        topMargin=72, bottomMargin=72, leftMargin=72, rightMargin=72
    )

    styles = getSampleStyleSheet()
    normal_style = ParagraphStyle(
        "Normal_CN", parent=styles["Normal"],
        fontName=font_name, fontSize=11, leading=18,
    )
    heading_style = ParagraphStyle(
        "Heading_CN", parent=styles["Heading1"],
        fontName=font_name, fontSize=16, leading=24,
        spaceAfter=12, spaceBefore=12,
    )

    elements = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            elements.append(Spacer(1, 12))
            continue
        # 转义 XML 特殊字符
        text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        style = heading_style if para.style.name.startswith("Heading") else normal_style
        try:
            elements.append(Paragraph(text, style))
        except Exception:
            elements.append(Paragraph(text, normal_style))

    if not elements:
        elements.append(Spacer(1, 1))

    pdf_doc.build(elements)
    return output_path


def word_to_text(input_path: Path, output_path: Path) -> Path:
    """Word → 纯文本"""
    from docx import Document
    doc = Document(str(input_path))
    lines = []
    for para in doc.paragraphs:
        lines.append(para.text)
    output_path.write_text("\n".join(lines), encoding="utf-8")
    return output_path


def word_to_html(input_path: Path, output_path: Path) -> Path:
    """Word → HTML"""
    from docx import Document
    import html as html_module

    doc = Document(str(input_path))
    body_parts = []

    for para in doc.paragraphs:
        text = html_module.escape(para.text)
        if not text.strip():
            body_parts.append("<br>")
            continue
        style_name = para.style.name.lower()
        if "heading 1" in style_name:
            body_parts.append(f"<h1>{text}</h1>")
        elif "heading 2" in style_name:
            body_parts.append(f"<h2>{text}</h2>")
        elif "heading 3" in style_name:
            body_parts.append(f"<h3>{text}</h3>")
        elif "list" in style_name:
            body_parts.append(f"<li>{text}</li>")
        else:
            body_parts.append(f"<p>{text}</p>")

    # 处理表格
    for table in doc.tables:
        body_parts.append('<table border="1" cellpadding="8" cellspacing="0" style="border-collapse:collapse;">')
        for i, row in enumerate(table.rows):
            cells = "".join(
                f'<{"th" if i == 0 else "td"} style="border:1px solid #ddd;padding:8px;">'
                f'{html_module.escape(cell.text)}</{"th" if i == 0 else "td"}>'
                for cell in row.cells
            )
            body_parts.append(f"<tr>{cells}</tr>")
        body_parts.append("</table>")

    html_content = f"""<!DOCTYPE html>
<html><head><meta charset="utf-8">
<title>{input_path.stem}</title>
<style>body{{font-family:system-ui,-apple-system,sans-serif;max-width:900px;margin:40px auto;padding:0 20px;line-height:1.8;color:#333;}}
table{{border-collapse:collapse;margin:20px 0;}}</style>
</head><body>
{''.join(body_parts)}
</body></html>"""
    output_path.write_text(html_content, encoding="utf-8")
    return output_path


def word_to_ppt(input_path: Path, output_path: Path) -> Path:
    """Word → PPT"""
    from docx import Document
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    doc = Document(str(input_path))
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 标题页布局
    blank_layout = prs.slide_layouts[6]  # blank layout

    current_slide = None
    current_text_frame = None
    y_position = Inches(1.5)
    title_done = False

    def add_new_slide():
        nonlocal current_slide, current_text_frame, y_position
        slide = prs.slides.add_slide(blank_layout)
        # 添加标题区域背景
        from pptx.util import Inches as In
        title_box = slide.shapes.add_textbox(In(0.8), In(0.4), In(11.5), In(1))
        tf = title_box.text_frame
        tf.word_wrap = True
        current_slide = slide
        current_text_frame = slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(11), Inches(5.5)
        ).text_frame
        current_text_frame.word_wrap = True
        y_position = Inches(0)
        return slide

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        style_name = para.style.name.lower()

        if "heading 1" in style_name or ("heading" not in style_name and not title_done):
            # 创建新幻灯片作为标题
            add_new_slide()
            title_tf = current_slide.shapes[-2].text_frame  # title box
            p = title_tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x1A, 0x56, 0xDB)
            title_done = True
            continue

        if "heading" in style_name:
            # 子标题 → 新幻灯片
            add_new_slide()
            title_tf = current_slide.shapes[-2].text_frame
            p = title_tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            continue

        if current_text_frame is None:
            add_new_slide()
            title_tf = current_slide.shapes[-2].text_frame
            p = title_tf.paragraphs[0]
            p.text = "内容"
            p.font.size = Pt(24)
            p.font.bold = True

        p = current_text_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.space_after = Pt(8)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    if not prs.slides:
        add_new_slide()
        title_tf = current_slide.shapes[-2].text_frame
        p = title_tf.paragraphs[0]
        p.text = input_path.stem
        p.font.size = Pt(28)
        p.font.bold = True

    prs.save(str(output_path))
    return output_path


def text_to_pdf(input_path: Path, output_path: Path) -> Path:
    """文本 → PDF"""
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    import os

    font_name = "Helvetica"
    for fp in ["C:/Windows/Fonts/simhei.ttf", "C:/Windows/Fonts/simsun.ttc"]:
        if os.path.exists(fp):
            try:
                pdfmetrics.registerFont(TTFont("ChineseFont", fp))
                font_name = "ChineseFont"
                break
            except Exception:
                continue

    text = input_path.read_text(encoding="utf-8", errors="replace")
    pdf_doc = SimpleDocTemplate(str(output_path), pagesize=A4)
    style = ParagraphStyle("Body", fontName=font_name, fontSize=11, leading=18)

    elements = []
    for line in text.split("\n"):
        line = line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        elements.append(Paragraph(line or "&nbsp;", style))

    pdf_doc.build(elements)
    return output_path


def text_to_word(input_path: Path, output_path: Path) -> Path:
    """文本 → Word"""
    from docx import Document
    text = input_path.read_text(encoding="utf-8", errors="replace")
    doc = Document()
    for line in text.split("\n"):
        doc.add_paragraph(line)
    doc.save(str(output_path))
    return output_path
